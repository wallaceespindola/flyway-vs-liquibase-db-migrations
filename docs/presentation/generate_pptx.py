#!/usr/bin/env python3
"""Generate flyway-vs-liquibase-deck.pptx from the content of flyway-vs-liquibase-slides.md.

Run:  python3 docs/presentation/generate_pptx.py

Every fact in this deck comes from the repository it lives in. Nothing is invented.

Author: Wallace Espindola <wallace.espindola@gmail.com>
        https://github.com/wallaceespindola/ | https://www.linkedin.com/in/wallaceespindola/
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

OUT = Path(__file__).resolve().parent / "flyway-vs-liquibase-deck.pptx"

# --- theme -----------------------------------------------------------------
BG = RGBColor(0x1E, 0x29, 0x3B)  # dark slate
PANEL = RGBColor(0x0F, 0x17, 0x2A)  # code / table body
PANEL_ALT = RGBColor(0x17, 0x22, 0x39)  # zebra row
TEXT = RGBColor(0xE2, 0xE8, 0xF0)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
ACCENT = RGBColor(0xFB, 0xBF, 0x24)  # amber, headings + rules
FLYWAY = RGBColor(0xF4, 0x3F, 0x5E)  # rose
LIQUIBASE = RGBColor(0x22, 0xD3, 0xEE)  # cyan
CODE_FG = RGBColor(0xCB, 0xD5, 0xE1)

BODY_FONT = "Calibri"
MONO_FONT = "Consolas"

W, H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = W - 2 * MARGIN


# --- primitives ------------------------------------------------------------
def new_slide(prs, notes):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.shapes.add_shape(1, 0, 0, W, H)  # 1 = rectangle
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    slide.notes_slide.notes_text_frame.text = notes.strip()
    return slide


def box(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, text, size, color, *, first=False, bold=False, font=BODY_FONT,
         space_after=6, bullet=False, align=PP_ALIGN.LEFT, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = ("•  " + text) if bullet else text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return p


def heading(slide, title, *, accent=ACCENT):
    tf = box(slide, MARGIN, Inches(0.42), CONTENT_W, Inches(0.75))
    para(tf, title, 30, TEXT, first=True, bold=True, space_after=0)
    rule = slide.shapes.add_shape(1, MARGIN, Inches(1.18), Inches(1.6), Emu(38100))
    rule.fill.solid()
    rule.fill.fore_color.rgb = accent
    rule.line.fill.background()
    rule.shadow.inherit = False


def bullets(slide, items, *, top=Inches(1.55), size=17, width=None, left=MARGIN):
    """items: list of str, or (str, color) tuples. Cap at ~6 for overflow safety."""
    tf = box(slide, left, top, width or CONTENT_W, H - top - Inches(0.5))
    for i, item in enumerate(items):
        text, color = item if isinstance(item, tuple) else (item, TEXT)
        para(tf, text, size, color, first=(i == 0), bullet=True, space_after=12)
    return tf


def code(slide, lines, *, top, left=MARGIN, width=None, size=12.5, caption=None,
         caption_color=ACCENT):
    width = width or CONTENT_W
    y = top
    if caption:
        tf = box(slide, left, y, width, Inches(0.28))
        para(tf, caption, 13, caption_color, first=True, bold=True, space_after=0)
        y += Inches(0.33)
    body = [ln for ln in lines.strip("\n").split("\n")]
    height = Inches(0.20) + Pt(size * 1.45) * len(body)
    panel = slide.shapes.add_shape(1, left, y, width, height)
    panel.fill.solid()
    panel.fill.fore_color.rgb = PANEL
    panel.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
    panel.line.width = Pt(0.75)
    panel.shadow.inherit = False
    tf = panel.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = Inches(0.16)
    tf.margin_top = tf.margin_bottom = Inches(0.10)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(body):
        colour = MUTED if line.lstrip().startswith(("--", "//", "#", "<!--")) else CODE_FG
        para(tf, line or " ", size, colour, first=(i == 0), font=MONO_FONT, space_after=0)
    return y + height


def table(slide, headers, rows, *, top=Inches(1.6), col_widths=None, size=12,
          header_size=12.5, left=MARGIN, width=None):
    width = width or CONTENT_W
    n_rows, n_cols = len(rows) + 1, len(headers)
    height = Inches(0.42) + Inches(0.36) * len(rows)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    tbl = shape.table

    if col_widths:
        total = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Emu(int(width * cw / total))

    for c, text in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x33, 0x41, 0x55)
        cell.margin_left = cell.margin_right = Inches(0.08)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf = cell.text_frame
        tf.word_wrap = True
        colour = FLYWAY if text == "Flyway" else LIQUIBASE if text == "Liquibase" else ACCENT
        para(tf, text, header_size, colour, first=True, bold=True, space_after=0)

    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if r % 2 else PANEL_ALT
            cell.margin_left = cell.margin_right = Inches(0.08)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            colour = TEXT
            bold = False
            if text == "Flyway":
                colour, bold = FLYWAY, True
            elif text == "Liquibase":
                colour, bold = LIQUIBASE, True
            elif text == "Tie":
                colour, bold = MUTED, True
            para(tf, text, size, colour, first=True, bold=bold, space_after=0)
    return tbl


def footer(slide, text="Flyway vs Liquibase  |  Wallace Espindola"):
    tf = box(slide, MARGIN, H - Inches(0.45), CONTENT_W, Inches(0.28))
    para(tf, text, 10, RGBColor(0x64, 0x74, 0x8B), first=True, space_after=0)


# --- slides ----------------------------------------------------------------
def slide_title(prs):
    s = new_slide(prs, """
This talk is backed by a running application, not a blog post. One Spring Boot process opens two
independent H2 databases. Flyway migrates one, Liquibase migrates the other, and the app reads both
back from INFORMATION_SCHEMA and diffs them at runtime.
Set the frame early: I am not selling either tool. Both work. The interesting question is which
trade-offs you are buying, and this deck answers that with code from the repo and one measured
result at the end.
""")
    band = s.shapes.add_shape(1, 0, Inches(2.35), W, Emu(57150))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    band.shadow.inherit = False

    tf = box(s, MARGIN, Inches(1.25), CONTENT_W, Inches(1.1))
    p = tf.paragraphs[0]
    p.space_after = Pt(0)
    for text, colour in (("Flyway", FLYWAY), ("  vs  ", MUTED), ("Liquibase", LIQUIBASE)):
        run = p.add_run()
        run.text = text
        run.font.size = Pt(54)
        run.font.bold = True
        run.font.color.rgb = colour
        run.font.name = BODY_FONT

    tf = box(s, MARGIN, Inches(2.65), CONTENT_W, Inches(2.6))
    para(tf, "A measured comparison, not an opinion", 24, TEXT, first=True, space_after=22)
    para(tf, "Wallace Espindola", 20, ACCENT, bold=True, space_after=4)
    para(tf, "Senior Software Engineer & Solution Architect", 15, MUTED, space_after=18)
    para(tf, "github.com/wallaceespindola   |   linkedin.com/in/wallaceespindola", 14, TEXT,
         space_after=6)
    para(tf, "github.com/wallaceespindola/flyway-vs-liquibase-db-migrations", 14, TEXT,
         space_after=18)
    para(tf, "Spring Boot 3.4.2  ·  Java 21  ·  Maven  ·  H2  ·  plain JDBC", 14, MUTED)


def slide_problem(prs):
    s = new_slide(prs, """
Everybody in the room has lived this. The symptom is not a dramatic outage, it is a slow
accumulation of uncertainty. You stop trusting that staging looks like production, so you stop
testing against staging, so defects reach production.
The root cause is that DDL is treated as an operational act rather than as source code. Two things
are missing: an ordered, reviewable record of every change, and a machine that applies exactly that
record to every environment.
""")
    heading(s, "The problem: schema drift and manual DDL")
    bullets(s, [
        "Application code is versioned, reviewed and rolled forward. The schema often is not.",
        "“Run this script on prod” in a chat message is not a deployment process.",
        "Environments diverge silently: dev has the column, prod has it with a different type.",
        "Nobody can answer “which changes has this database actually seen?”",
        "Rebuilding an environment from scratch becomes archaeology.",
    ])
    footer(s)


def slide_why(prs):
    s = new_slide(prs, """
Frame these four as the baseline, not the differentiator. Any tool that fails one of these is not a
candidate.
The checksum point deserves a beat: both tools hash applied changes, so editing a migration that has
already run fails validation instead of silently drifting. That single behaviour prevents a large
class of "it works on my machine" incidents.
The automation point matters for Spring Boot: both engines run before JPA initialises, so the schema
is guaranteed correct before the first entity is touched.
""")
    heading(s, "Why migration tooling")
    tf = box(s, MARGIN, Inches(1.5), CONTENT_W, Inches(0.4))
    para(tf, "Four properties that scripts in a wiki do not give you:", 16, MUTED, first=True)
    bullets(s, [
        "Ordering — changes apply in a defined sequence, once, everywhere.",
        "Bookkeeping — a table inside the database records what has been applied.",
        "Integrity — a checksum detects edits to already-applied changes.",
        "Automation — migration runs at startup or in the pipeline, not by hand.",
    ], top=Inches(2.1), size=18)
    tf = box(s, MARGIN, Inches(5.1), CONTENT_W, Inches(0.5))
    para(tf, "Flyway and Liquibase both deliver all four. Everything else is about how.",
         18, ACCENT, first=True, bold=True)
    footer(s)


def slide_how_flyway(prs):
    s = new_slide(prs, """
This is the whole model. There is no changelog file, no registry, no include list. Convention over
configuration in the most literal sense: discovery is a directory scan and an ordering rule.
Point at the bootstrap code - it is from config/FlywayConfig.java in the repo, unedited. Five
builder calls and Flyway is live. That brevity is Flyway's core value proposition and it is
genuinely hard to argue with.
validateOnMigrate(true) is the checksum enforcement. cleanDisabled(true) in the repo blocks the
destructive clean command, which you always want switched off outside a laptop.
""")
    heading(s, "How Flyway works", accent=FLYWAY)
    bullets(s, [
        "Drop SQL files into a location: V1__create_category_table.sql",
        "The filename IS the configuration: V = versioned, 1 = order, rest = description.",
        "R__ prefix marks a repeatable migration, re-applied when its checksum changes.",
        "Flyway scans, sorts by version, runs what is not yet in flyway_schema_history.",
        "Scripts are raw SQL in the target dialect. No abstraction layer.",
    ], top=Inches(1.5), size=16)
    code(s, """
Flyway.configure()
      .dataSource(dataSource)
      .locations("classpath:db/migration")
      .baselineOnMigrate(true)
      .validateOnMigrate(true)
      .load();
""", top=Inches(4.35), caption="config/FlywayConfig.java", caption_color=FLYWAY)
    footer(s)


def slide_how_liquibase(prs):
    s = new_slide(prs, """
Liquibase needs more knobs than Flyway, and that is not an accident: it models changes abstractly,
so it needs to be told what to include, in what context, and against which schema.
The abstraction is the product. When you write createTable instead of CREATE TABLE, Liquibase
decides what H2, PostgreSQL or Oracle should actually receive. That is what unlocks portability, and
it is also what makes the diff harder to read for a reviewer who only knows SQL.
This bootstrap is from config/LiquibaseConfig.java, unedited. SpringLiquibase implements
InitializingBean, so the changelog is applied during bean initialisation - the same lifecycle
position as Flyway's migrate().
""")
    heading(s, "How Liquibase works", accent=LIQUIBASE)
    bullets(s, [
        "A changeset is the unit of change: id + author + source file identify it.",
        "Changesets live in XML, YAML, JSON or SQL — all first-class, all mixable.",
        "A master changelog lists every included file. Ordering is declared, not inferred.",
        "Changesets are database-agnostic; Liquibase emits the dialect at runtime.",
        "preConditions, context, labels, rollback, runOnChange are per-changeset attributes.",
    ], top=Inches(1.5), size=16)
    code(s, """
SpringLiquibase liquibase = new SpringLiquibase();
liquibase.setDataSource(dataSource);
liquibase.setChangeLog("classpath:db/changelog/db.changelog-master.yaml");
liquibase.setContexts("demo");
liquibase.setDefaultSchema("PUBLIC");
""", top=Inches(4.5), caption="config/LiquibaseConfig.java", caption_color=LIQUIBASE)
    footer(s)


def slide_experiment(prs):
    s = new_slide(prs, """
The design matters because it removes the usual hand-waving. Both engines are explicitly wired, with
Spring Boot's DataSourceAutoConfiguration excluded, so nothing is hidden behind auto-configuration
magic. You see the real bootstrap of each tool.
Each JdbcTemplate carries @DependsOn on its migration bean. That guarantees no read can happen
before the migrations have run, so the comparison always reflects a fully migrated database.
Then SchemaInspectionService queries INFORMATION_SCHEMA on both sides - it never trusts the
migration scripts, it reads what the database actually contains.
""")
    heading(s, "The experiment design")
    tf = box(s, MARGIN, Inches(1.45), CONTENT_W, Inches(0.4))
    para(tf, "Two independent H2 databases. One identical logical schema. One JVM, port 8080.",
         17, MUTED, first=True)
    table(s, ["", "Flyway side", "Liquibase side"], [
        ["Database", "./data/flywaydb", "./data/liquibasedb"],
        ["Wiring", "FlywayConfig @Configuration", "LiquibaseConfig @Configuration"],
        ["Runner", "Flyway.migrate() on bean init", "SpringLiquibase on bean init"],
        ["Reader", "flywayJdbcTemplate (@DependsOn)", "liquibaseJdbcTemplate (@DependsOn)"],
        ["Changes", "6 versioned + repeatable SQL files", "7 changesets in 6 changelog files"],
    ], top=Inches(2.05), col_widths=[1.1, 2.6, 2.6], size=13.5)
    tf = box(s, MARGIN, Inches(4.75), CONTENT_W, Inches(0.5))
    para(tf, "The app then reads both schemas from INFORMATION_SCHEMA and diffs them at runtime.",
         17, ACCENT, first=True, bold=True)
    footer(s)


def slide_flyway_walkthrough(prs):
    s = new_slide(prs, """
Walk the table left to right. Note that V3 seeds data through a versioned migration - the seed
becomes part of schema history and applies exactly once per environment, in order. That is a
deliberate choice and it is the right one for reference data.
V5 is the classic expand step of an expand/contract rollout: add a column with a default, backfill
it, index it. All in one migration because it is one logical change.
The R__ file has no version number. Flyway re-runs it whenever the file's checksum changes, which is
exactly what you want for views and stored procedures - objects you redefine rather than version.
""")
    heading(s, "Flyway migrations: the walkthrough", accent=FLYWAY)
    tf = box(s, MARGIN, Inches(1.45), CONTENT_W, Inches(0.35))
    para(tf, "src/main/resources/db/migration", 15, MUTED, first=True, font=MONO_FONT)
    table(s, ["File", "What it does"], [
        ["V1__create_category_table.sql", "category table + unique name + index"],
        ["V2__create_product_table.sql", "product + FK, 2 check constraints, 2 indexes"],
        ["V3__seed_reference_data.sql", "3 categories, 5 products"],
        ["V4__add_product_audit_table.sql", "product_audit + FK + backfill from product"],
        ["V5__add_product_active_flag.sql", "ADD COLUMN active + backfill + index"],
        ["R__product_catalog_view.sql", "v_product_catalog, repeatable"],
    ], top=Inches(1.9), col_widths=[2.1, 3.0], size=13.5)
    tf = box(s, MARGIN, Inches(5.05), CONTENT_W, Inches(0.5))
    para(tf, "6 migrations applied.", 20, FLYWAY, first=True, bold=True)
    footer(s)


def slide_flyway_code(prs):
    s = new_slide(prs, """
This is the ergonomic argument for Flyway in one slide. Anyone who reads SQL can review these two
files with no additional vocabulary. There is nothing between the intent and the statement.
The cost is written on the same slide: this is H2 dialect. Move to PostgreSQL and BOOLEAN NOT NULL
DEFAULT TRUE is fine, but AUTO_INCREMENT in V1 is not. You would rewrite the files.
Note also what is missing from V5: any way to undo it. Flyway Community has no undo. Reverting means
writing V6__drop_product_active_flag.sql. We come back to this on the rollback slide.
""")
    heading(s, "Flyway code: expand/contract and repeatable", accent=FLYWAY)
    y = code(s, """
ALTER TABLE product
    ADD COLUMN active BOOLEAN NOT NULL DEFAULT TRUE;

UPDATE product SET active = FALSE WHERE stock_quantity = 0;

CREATE INDEX idx_product_active ON product (active);
""", top=Inches(1.5), caption="V5__add_product_active_flag.sql", caption_color=FLYWAY)
    code(s, """
CREATE OR REPLACE VIEW v_product_catalog AS
SELECT p.id AS product_id, p.sku, p.name AS product_name,
       p.price, p.stock_quantity, p.active,
       c.id AS category_id, c.name AS category_name
FROM product p JOIN category c ON c.id = p.category_id;
""", top=y + Inches(0.35),
         caption="R__product_catalog_view.sql  —  re-runs when the checksum changes",
         caption_color=FLYWAY)
    footer(s)


def slide_liquibase_walkthrough(prs):
    s = new_slide(prs, """
The format mixing is not showing off - it is how real Liquibase projects look. Structural changes go
in XML or YAML for portability, and you drop to raw SQL only where a vendor feature has no portable
tag. 002 does exactly that: check constraints have no Liquibase tag, so there is an inline sql block.
Count carefully: six files, seven changesets, because 005 splits the schema change from the data
backfill. That split is itself the point - they carry different labels, schema-evolution and
data-backfill, so they can be selected independently at runtime.
Compare with the Flyway side: six migrations there, seven changesets here, same resulting schema.
""")
    heading(s, "Liquibase changelog: the walkthrough", accent=LIQUIBASE)
    tf = box(s, MARGIN, Inches(1.42), CONTENT_W, Inches(0.35))
    para(tf, "db/changelog/db.changelog-master.yaml includes six files, in three formats", 15,
         MUTED, first=True)
    table(s, ["File", "Format", "What it demonstrates"], [
        ["001-create-category-table.xml", "XML", "portable createTable + explicit rollback"],
        ["002-create-product-table.yaml", "YAML", "terser diff; raw SQL for check constraints"],
        ["003-seed-reference-data.sql", "SQL", "--liquibase formatted sql, --rollback directives"],
        ["004-add-product-audit-table.xml", "XML", "preConditions + rollback"],
        ["005-add-product-active-flag.xml", "XML", "context + labels, two changesets"],
        ["006-product-catalog-view.xml", "XML", "runOnChange=\"true\""],
    ], top=Inches(1.88), col_widths=[2.2, 0.7, 3.2], size=13)
    tf = box(s, MARGIN, Inches(5.0), CONTENT_W, Inches(0.6))
    para(tf, "7 changesets applied — 005 contains two: the column and the backfill.",
         20, LIQUIBASE, first=True, bold=True)
    footer(s)


def slide_liquibase_code(prs):
    s = new_slide(prs, """
Two capabilities are on display and neither exists in Flyway Community.
First, preConditions. If the product table is missing, the changeset is marked as run instead of
exploding halfway through a deploy. onFail has several modes - MARK_RAN, CONTINUE, HALT, WARN - so
you choose the failure semantics per changeset. Flyway has no declarative equivalent; you would
write defensive SQL or a Java migration.
Second, rollback. This block makes "liquibase rollbackCount 1" a supported, testable operation.
Liquibase can infer rollbacks for most structural changes, but declaring them explicitly keeps
intent visible and survives future refactoring of the forward change.
""")
    heading(s, "Liquibase code: preconditions and rollback", accent=LIQUIBASE)
    code(s, """
<changeSet id="004-add-product-audit-table" author="wallaceespindola" context="demo">
  <preConditions onFail="MARK_RAN"
                 onFailMessage="product table missing, skipping audit table">
    <tableExists tableName="product"/>
  </preConditions>

  <createTable tableName="product_audit"> ... </createTable>
  <createIndex tableName="product_audit" indexName="idx_product_audit_product"> ... </createIndex>

  <rollback>
    <dropIndex tableName="product_audit" indexName="idx_product_audit_product"/>
    <dropTable tableName="product_audit"/>
  </rollback>
</changeSet>
""", top=Inches(1.5), size=12.5,
         caption="db/changelog/changes/004-add-product-audit-table.xml", caption_color=LIQUIBASE)
    tf = box(s, MARGIN, Inches(6.15), CONTENT_W, Inches(0.5))
    para(tf, "This changeset has no Flyway Community equivalent.", 18, ACCENT, first=True,
         bold=True)
    footer(s)


def slide_side_by_side(prs):
    s = new_slide(prs, """
Same table, two philosophies, and the trade-off is visible in the character count.
The Flyway version is shorter and every reviewer already speaks it. The Liquibase version is longer
and requires knowing what the tags emit - but it names its constraints explicitly and it is not
bound to H2. AUTO_INCREMENT in the SQL version is H2 and MySQL syntax; PostgreSQL wants GENERATED
ALWAYS AS IDENTITY. Liquibase's autoIncrement="true" produces the right thing on either.
This is the whole comparison in miniature: directness versus abstraction. Neither is free.
""")
    heading(s, "Side by side: the same table, two ways")
    y = code(s, """
CREATE TABLE category (
    id          BIGINT AUTO_INCREMENT PRIMARY KEY,
    name        VARCHAR(100) NOT NULL,
    created_at  TIMESTAMP NOT NULL DEFAULT CURRENT_TIMESTAMP,
    CONSTRAINT uk_category_name UNIQUE (name));
""", top=Inches(1.5), size=12.5,
         caption="Flyway — V1__create_category_table.sql", caption_color=FLYWAY)
    code(s, """
<createTable tableName="category">
  <column name="id" type="BIGINT" autoIncrement="true">
    <constraints primaryKey="true" primaryKeyName="pk_category" nullable="false"/></column>
  <column name="name" type="VARCHAR(100)">
    <constraints nullable="false" unique="true" uniqueConstraintName="uk_category_name"/></column>
</createTable>
""", top=y + Inches(0.35), size=12.5,
         caption="Liquibase — 001-create-category-table.xml", caption_color=LIQUIBASE)
    footer(s)


def slide_bookkeeping(prs):
    s = new_slide(prs, """
This slide is measured, not editorial - it is what the two tables actually contain in the running
demo. Two rows deserve emphasis in opposite directions.
Flyway records execution time per migration in milliseconds. Liquibase does not persist a
per-changeset duration at all. If you want to know which migration is making your deploys slow,
Flyway tells you for free.
Liquibase records author, contexts, labels and deployment id. Flyway's attribution lives only in git
history. If auditors want to see who changed the schema without leaving the database, Liquibase
answers that and Flyway does not.
Also note Liquibase needs a second table for locking, Flyway uses a database-level lock.
""")
    heading(s, "Bookkeeping: what each engine records")
    table(s, ["", "flyway_schema_history", "DATABASECHANGELOG"], [
        ["Identity", "version (single ordered namespace)", "id + author + filename"],
        ["Author", "not recorded", "mandatory attribute"],
        ["Checksum", "CRC32", "MD5"],
        ["Execution time", "yes, milliseconds", "not recorded"],
        ["Contexts / labels", "not applicable", "recorded per changeset"],
        ["Deployment id", "not recorded", "recorded"],
        ["Locking", "database-level lock during run", "DATABASECHANGELOGLOCK table"],
    ], top=Inches(1.6), col_widths=[1.3, 2.5, 2.5], size=13.5)
    footer(s)


def slide_history_api(prs):
    s = new_slide(prs, """
This is the finding that surprised me most when building the demo, and it is a real operational
difference if you embed the engine in an application.
flyway.info() hands you objects. Applied and pending, already ordered, with state and checksum. Zero
SQL. That is FlywayHistoryService in the repo - the entire class is a stream over that array.
On the Liquibase side, LiquibaseHistoryService is a JdbcTemplate query against DATABASECHANGELOG,
because there is no equivalent. It works, and the table gives you more columns than Flyway's does,
but you are now coupled to Liquibase's internal table shape.
Also worth saying honestly: Liquibase only records what it has already run. Discovering pending
changesets requires a full changelog parse.
""")
    heading(s, "Reading the history: two different APIs")
    y = code(s, """
MigrationInfo[] all = flyway.info().all();
// applied + pending, ordered, with state, checksum and executionTime attached
""", top=Inches(1.5), caption="Flyway — an embedded status API, no SQL needed",
         caption_color=FLYWAY)
    code(s, """
SELECT ID, AUTHOR, FILENAME, DATEEXECUTED, ORDEREXECUTED, EXECTYPE,
       MD5SUM, DESCRIPTION, COMMENTS, CONTEXTS, LABELS, DEPLOYMENT_ID
FROM DATABASECHANGELOG ORDER BY ORDEREXECUTED
""", top=y + Inches(0.4), caption="Liquibase — query the bookkeeping table yourself",
         caption_color=LIQUIBASE)
    tf = box(s, MARGIN, Inches(5.3), CONTENT_W, Inches(0.6))
    para(tf, "Liquibase has no lightweight read-only status API. The table is the public interface.",
         17, ACCENT, first=True, bold=True)
    footer(s)


def slide_result(prs):
    s = new_slide(prs, """
This is the headline. Six Flyway migrations and seven Liquibase changesets converge on identical
business schemas. Tables, views and columns with their data types all match.
Be precise about what is compared and what is not. ComparisonService diffs tables, views and
columns. It deliberately excludes indexes, because H2 auto-generates constraint-backing indexes
under generated names that legitimately differ between the two engines - including them would
report noise as drift. That exclusion is documented in the code, not hidden.
The conclusion to draw: the choice between these tools is not about what schema you end up with. It
is entirely about the process around getting there.
""")
    heading(s, "The measured result")
    tf = box(s, MARGIN, Inches(1.42), CONTENT_W, Inches(0.35))
    para(tf, "GET /api/v1/comparison on the running application returns:", 16, MUTED, first=True)
    code(s, """
"schemasEquivalent": true,
"schemaDifferences": []
""", top=Inches(1.88), size=17)
    bullets(s, [
        "Tables on both sides: category, product, product_audit",
        "View on both sides: v_product_catalog",
        "Columns compared as TABLE.COLUMN:TYPE — zero differences",
        "Only difference: flyway_schema_history vs DATABASECHANGELOG + ...LOCK",
    ], top=Inches(3.6), size=17)
    footer(s)


def slide_matrix_1(prs):
    s = new_slide(prs, """
This matrix is served live at /api/v1/comparison/features and lives in FeatureMatrix.java. Eighteen
rows total, split across three slides. The "Edge" column is my judgement, stated as such -
everything to its left is factual.
Notice the pattern already forming: Flyway wins the rows about humans reading and writing changes,
Liquibase wins the rows about machines executing them in varied environments.
Do not linger on every row. Read two or three and move on; the deck is a reference the audience can
re-read.
""")
    heading(s, "Feature matrix (1 of 3) — authoring")
    table(s, ["Capability", "Flyway", "Liquibase", "Edge"], [
        ["Change format", "Plain SQL only (Community); Java migrations",
         "XML, YAML, JSON or SQL, mixable", "Liquibase"],
        ["Learning curve", "One naming convention: V1__name.sql",
         "Changeset model, changelog composition, tags", "Flyway"],
        ["Database portability", "None — scripts are in the target dialect",
         "Abstract changesets; dialect at runtime", "Liquibase"],
        ["Migration discovery", "Convention: scan a location, order by version",
         "Explicit master changelog lists every file", "Flyway"],
        ["Repeatable changes", "R__ scripts, re-run on checksum change",
         "runOnChange=\"true\" on any changeset", "Tie"],
        ["Review ergonomics", "Diffs are SQL — every reviewer reads it",
         "Diffs are tags — must know what they emit", "Flyway"],
    ], top=Inches(1.6), col_widths=[1.3, 2.4, 2.4, 0.75], size=12)
    footer(s)


def slide_matrix_2(prs):
    s = new_slide(prs, """
This is where Liquibase collects most of its wins, and they are substantive rather than cosmetic.
Rollback, conditional execution and diff are three genuinely different capabilities, not three names
for the same one.
The Spring Boot row is worth calling out as a tie because people assume otherwise. Both engines are
auto-configured, both run before JPA initialises, both are one property block away from working. Our
demo bypasses auto-configuration on purpose, to show the real bootstrap.
The timing row is the one place Flyway is ahead on observability, and it is genuinely useful when a
deploy window is tight.
""")
    heading(s, "Feature matrix (2 of 3) — execution")
    table(s, ["Capability", "Flyway", "Liquibase", "Edge"], [
        ["Rollback / undo", "Not in Community; undo is a Teams feature",
         "Built in: inferred, or declared <rollback>", "Liquibase"],
        ["Conditional execution", "Placeholders, per-environment locations",
         "Preconditions, contexts, labels per changeset", "Liquibase"],
        ["Drift detection / diff", "Not available in Community",
         "diff and diffChangeLog generate the delta", "Liquibase"],
        ["Concurrency safety", "Database-level lock for the migration run",
         "Dedicated DATABASECHANGELOGLOCK table", "Tie"],
        ["Spring Boot integration", "spring.flyway.*, runs before JPA",
         "spring.liquibase.*, same lifecycle position", "Tie"],
        ["Execution timing", "Yes, per migration, in milliseconds",
         "No per-changeset duration persisted", "Flyway"],
    ], top=Inches(1.6), col_widths=[1.4, 2.35, 2.35, 0.75], size=12)
    footer(s)


def slide_matrix_3(prs):
    s = new_slide(prs, """
Two rows here decide real procurement conversations.
The licensing row: the capabilities most teams eventually want - undo and drift detection - are paid
features in Flyway and open source in Liquibase. That is not a knock on Flyway's business model, but
it belongs in your evaluation, because "Flyway Community is enough" often stops being true about
eighteen months in.
The merge conflict row is the one experienced teams nod at. Two branches both adding V6 collide
immediately and obviously. Two branches both appending to a master changelog produce a conflict git
will happily auto-resolve into the wrong order, and nothing fails until deploy.
""")
    heading(s, "Feature matrix (3 of 3) — operations")
    table(s, ["Capability", "Flyway", "Liquibase", "Edge"], [
        ["History bookkeeping", "version, description, checksum, timing",
         "id, author, contexts, labels, deployment id", "Liquibase"],
        ["Embedded status API", "flyway.info() returns applied and pending",
         "No read API — query DATABASECHANGELOG", "Flyway"],
        ["Authorship tracking", "Not recorded; lives in version control",
         "author is mandatory on every changeset", "Liquibase"],
        ["Checksum protection", "CRC32; editing an applied migration fails",
         "MD5; same, plus runOnChange opt-out", "Tie"],
        ["Merge conflict profile", "Same version in two branches collides loudly",
         "Conflicts in the master include list auto-merge wrongly", "Flyway"],
        ["Licensing of advanced", "undo, dry-run, drift are paid Teams/Enterprise",
         "Rollback and diff open source; policy checks Pro", "Liquibase"],
    ], top=Inches(1.6), col_widths=[1.4, 2.35, 2.35, 0.75], size=12)
    footer(s)


def slide_rollback(prs):
    s = new_slide(prs, """
Be honest about how much this matters in practice, because the honest answer is "less than the
marketing suggests, but not zero".
Most production incidents are not fixed by rolling the schema back. Once data has been written
against the new shape, an automated rollback destroys it. The disciplined pattern - expand, migrate,
contract - makes forward-only recovery viable, and plenty of high-functioning teams run Flyway
Community forever without missing undo.
Where rollback genuinely pays: pre-production. Tearing a test environment back to a known point,
rehearsing a release, iterating on a changeset locally. Running rollbackCount 1 instead of
rebuilding the database is a real productivity gain, and it is free in Liquibase.
""")
    heading(s, "The rollback story")
    y = code(s, """
-- V6__drop_product_audit_table.sql   (the only Community option)
DROP TABLE product_audit;
""", top=Inches(1.5),
         caption="Flyway Community — no undo; reverting means a new forward migration",
         caption_color=FLYWAY)
    code(s, """
<rollback>
  <dropIndex tableName="product_audit" indexName="idx_product_audit_product"/>
  <dropTable tableName="product_audit"/>
</rollback>
""", top=y + Inches(0.4),
         caption="Liquibase — rollback declared next to the change it reverses",
         caption_color=LIQUIBASE)
    tf = box(s, MARGIN, Inches(5.1), CONTENT_W, Inches(0.9))
    para(tf, "liquibase rollbackCount 1 is a supported operation.", 18, ACCENT, first=True,
         bold=True, space_after=6)
    para(tf, "Every changeset in this repository declares one.", 16, MUTED)
    footer(s)


def slide_portability(prs):
    s = new_slide(prs, """
The nuance is the point of this slide. Liquibase's portability is real, and it is bounded.
Look at 002-create-product-table.yaml in the repo. The table, columns and indexes are portable tags.
The two check constraints are an inline sql block, with a comment saying exactly why: there is no
portable tag for them. So that file is portable in part and dialect-bound in part.
Then ask the room the question that actually settles it: how many times has your team changed
database engine? For most product teams the answer is zero, and portability is paying rent it does
not earn. For a vendor shipping the same product onto customer-chosen databases, it is the entire
reason to pick Liquibase.
""")
    heading(s, "The portability story")
    bullets(s, [
        "Flyway scripts are the target dialect. AUTO_INCREMENT is H2 and MySQL;",
        "PostgreSQL wants GENERATED ALWAYS AS IDENTITY. Changing engine means rewriting.",
        "Liquibase changesets are abstract: autoIncrement=\"true\" becomes what the target needs.",
        "But portability holds only while you stay inside the tag vocabulary.",
        "The repo shows the seam: check constraints have no portable tag, so 002 drops to <sql>.",
    ], top=Inches(1.55), size=17)
    code(s, """
- sql:
    comment: Check constraints have no portable Liquibase tag, so raw SQL is used here
    sql: >-
      ALTER TABLE product ADD CONSTRAINT ck_product_price_positive CHECK (price >= 0);
""", top=Inches(4.6), size=12,
         caption="002-create-product-table.yaml — portability has a boundary",
         caption_color=LIQUIBASE)
    footer(s)


def slide_review(prs):
    s = new_slide(prs, """
This is the slide that changes minds in rooms full of practitioners, because it is about the daily
cost rather than the feature list.
The Flyway conflict is the good kind of failure: it happens at merge time, it is obvious, and the
fix is renaming a file. Some teams add a CI check that version numbers are unique - trivially cheap.
The Liquibase conflict is the bad kind: syntactically valid, semantically wrong, discovered later.
The mitigations are real - one changelog file per release, or directory-based inclusion with
includeAll - but they are conventions your team has to adopt and enforce, not defaults.
If your team is large and merges often, weigh this row heavily.
""")
    heading(s, "Review and merge-conflict ergonomics")
    half = (CONTENT_W - Inches(0.5)) / 2
    tf = box(s, MARGIN, Inches(1.5), half, Inches(0.4))
    para(tf, "Flyway", 22, FLYWAY, first=True, bold=True)
    bullets(s, [
        "The pull request diff is SQL. Every reviewer already reads it.",
        "Two branches adding V6__ collide on the filename.",
        "Loud, immediate, unmissable — the fix is renaming a file.",
    ], top=Inches(2.1), size=15, left=MARGIN, width=half)

    tf = box(s, MARGIN + half + Inches(0.5), Inches(1.5), half, Inches(0.4))
    para(tf, "Liquibase", 22, LIQUIBASE, first=True, bold=True)
    bullets(s, [
        "The diff is XML or YAML tags; reviewers must know what each emits.",
        "Two branches appending to the master changelog conflict in the include list.",
        "Git can auto-merge that list into the wrong order. Nothing fails until deploy.",
    ], top=Inches(2.1), size=15, left=MARGIN + half + Inches(0.5), width=half)
    footer(s)


def slide_decision(prs):
    s = new_slide(prs, """
This is the slide people photograph. Give it time.
Reframe the decision so nobody leaves thinking one tool is better. Flyway optimises for the common
case: one database, SQL-literate team, forward-only releases. Liquibase optimises for variability:
many databases, many environments, changes that must be conditional or reversible.
The most common mistake is picking Liquibase for portability you will never exercise, and paying the
review-ergonomics tax every single day for it. The second most common is picking Flyway Community
and discovering eighteen months later that undo and drift detection are behind a licence.
Whichever you choose: pick one, use it for everything, and never apply DDL by hand again. That
decision matters more than which name you pick.
""")
    heading(s, "Decision guide")
    half = (CONTENT_W - Inches(0.5)) / 2
    tf = box(s, MARGIN, Inches(1.5), half, Inches(0.4))
    para(tf, "Choose Flyway when", 22, FLYWAY, first=True, bold=True)
    bullets(s, [
        "You target one database and expect to keep targeting it.",
        "Your team is fluent in SQL and wants readable review diffs.",
        "Forward-only migration fits your release process.",
        "You want the smallest thing between intent and executed statement.",
    ], top=Inches(2.1), size=15, left=MARGIN, width=half)

    tf = box(s, MARGIN + half + Inches(0.5), Inches(1.5), half, Inches(0.4))
    para(tf, "Choose Liquibase when", 22, LIQUIBASE, first=True, bold=True)
    bullets(s, [
        "You ship to several engines, or customers choose the engine.",
        "You need rollback, preconditions or conditional execution first-class.",
        "Audit wants author, context and deployment id inside the database.",
        "You want drift detection and diff without a commercial licence.",
    ], top=Inches(2.1), size=15, left=MARGIN + half + Inches(0.5), width=half)
    footer(s)


def slide_architecture(prs):
    s = new_slide(prs, """
Trace the flow once, top to bottom. Two configs, two datasources, two migration runners, two
JdbcTemplates guarded by @DependsOn so nothing can read before migration completes.
Both history services implement the same MigrationHistoryProvider interface, which is what lets
ComparisonService treat the engines symmetrically. The implementations are asymmetric for a reason
that is itself a finding: one calls an API, the other writes SQL.
FeatureMatrix is a static immutable list, not a database table. It is documentation that happens to
be served over HTTP, and it changes only when the tools do.
Excluding DataSourceAutoConfiguration is deliberate: with two datasources, explicit wiring is both
correct and more instructive.
""")
    heading(s, "Architecture of the demo app")
    code(s, """
DbMigrationComparisonApplication   (Spring Boot 3.4.2, Java 21, port 8080)
|
+-- FlywayConfig      -> flywayDataSource    -> Flyway.migrate()  -> flywayJdbcTemplate
+-- LiquibaseConfig   -> liquibaseDataSource -> SpringLiquibase   -> liquibaseJdbcTemplate
|
+-- FlywayHistoryService     (flyway.info())                  --+
+-- LiquibaseHistoryService  (SELECT ... DATABASECHANGELOG)     +--> ComparisonService
+-- SchemaInspectionService  (INFORMATION_SCHEMA x 2)         --+       -> ComparisonReport
+-- FeatureMatrix            (18 editorial rows)
""", top=Inches(1.6), size=13)
    tf = box(s, MARGIN, Inches(4.8), CONTENT_W, Inches(0.6))
    para(tf, "DataSourceAutoConfiguration is excluded on purpose — both engines are wired "
             "explicitly.", 17, ACCENT, first=True, bold=True)
    footer(s)


def slide_run(prs):
    s = new_slide(prs, """
If you are demoing live, this is the moment. Start with /api/v1/comparison and scroll to
schemasEquivalent: true - that is the payoff of the whole talk in one field.
Then open /api/v1/migrations and put the two arrays side by side. The Flyway entries carry
executionTimeMs and author "n/a". The Liquibase entries carry a real author and a null
executionTimeMs. That contrast in raw JSON makes the bookkeeping slide concrete.
If you have the H2 console open, connect to jdbc:h2:file:./data/flywaydb and then ./data/liquibasedb
and show the two bookkeeping tables next to each other. Credentials are sa with an empty password.
Everything runs on a laptop with no Docker and no external database.
""")
    heading(s, "How to run it")
    code(s, """
git clone https://github.com/wallaceespindola/flyway-vs-liquibase-db-migrations
cd flyway-vs-liquibase-db-migrations
mvn spring-boot:run
""", top=Inches(1.5), size=14)
    table(s, ["Endpoint", "Shows"], [
        ["/api/v1/comparison", "full side-by-side report + schemasEquivalent"],
        ["/api/v1/comparison/features", "the 18-row feature matrix"],
        ["/api/v1/migrations", "both engines' status"],
        ["/api/v1/migrations/{engine}/schema", "one engine's schema snapshot"],
        ["/api/v1/catalog/{engine}", "seeded data through v_product_catalog"],
        ["/swagger-ui.html  ·  /h2-console  ·  /api/v1/health", "docs, DB console, health"],
    ], top=Inches(3.25), col_widths=[2.2, 3.0], size=13)
    footer(s)


def slide_questions(prs):
    s = new_slide(prs, """
Close on the one-line summary: both tools produce the same schema, so choose on process, not on
outcome.
Questions you should expect:
- "Can I migrate from Flyway to Liquibase?" Yes, in both directions. Liquibase can generate a
  changelog from an existing database, and you mark everything up to now as already-run.
- "What about using both?" Technically possible with separate schemas, but you now maintain two
  mental models. Do not.
- "Does this work outside Spring Boot?" Both have Maven, Gradle and CLI interfaces. Nothing in this
  comparison depends on Spring.
- "Which do you use?" Answer honestly, and say why in terms of the decision-guide slide.
""")
    band = s.shapes.add_shape(1, 0, Inches(2.5), W, Emu(57150))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    band.shadow.inherit = False

    tf = box(s, MARGIN, Inches(1.5), CONTENT_W, Inches(1.0))
    para(tf, "Questions", 48, TEXT, first=True, bold=True, space_after=0)

    tf = box(s, MARGIN, Inches(2.85), CONTENT_W, Inches(3.0))
    para(tf, "Wallace Espindola", 22, ACCENT, first=True, bold=True, space_after=4)
    para(tf, "Senior Software Engineer & Solution Architect", 15, MUTED, space_after=20)
    para(tf, "github.com/wallaceespindola/flyway-vs-liquibase-db-migrations", 16, TEXT,
         space_after=6)
    para(tf, "github.com/wallaceespindola   |   linkedin.com/in/wallaceespindola", 16, TEXT,
         space_after=24)
    para(tf, "Both tools produce the same schema. Choose on process, not on outcome.",
         18, TEXT, italic=True)
    footer(s)


SLIDES = [
    slide_title, slide_problem, slide_why, slide_how_flyway, slide_how_liquibase,
    slide_experiment, slide_flyway_walkthrough, slide_flyway_code,
    slide_liquibase_walkthrough, slide_liquibase_code, slide_side_by_side,
    slide_bookkeeping, slide_history_api, slide_result,
    slide_matrix_1, slide_matrix_2, slide_matrix_3,
    slide_rollback, slide_portability, slide_review, slide_decision,
    slide_architecture, slide_run, slide_questions,
]


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    for build in SLIDES:
        build(prs)
    prs.save(OUT)

    # self-check: reopen and assert the deck is real
    check = Presentation(OUT)
    assert len(check.slides) == len(SLIDES), "slide count mismatch"
    assert all(s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip()
               for s in check.slides), "a slide is missing speaker notes"
    print(f"{OUT}  {len(check.slides)} slides  {OUT.stat().st_size:,} bytes")


if __name__ == "__main__":
    main()
